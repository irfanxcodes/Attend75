"""
Notes Ingestion Service

Processes uploaded notes files (PDF/PPTX/DOCX) to extract structured
problems and step-by-step solutions. Pipeline:

  1. Extract text (typed → direct; scanned PDF → Tesseract → Gemini Vision fallback)
  2. Phase 1 LLM call: extract problem list (question, topic, answer)
  3. Phase 2 LLM calls (per problem): pre-compute adjustments, then generate
     text steps + one table per accounting statement
  4. Persist to notes_problem_sets / notes_problems / notes_solution_steps
  5. Delete uploaded file from disk

Called as a FastAPI BackgroundTask — must be self-contained.
"""

import base64
import io
import json
import logging
import os
import time
import uuid
from datetime import datetime
from pathlib import Path
from statistics import mean
from typing import Optional, Literal

import requests
from pydantic import BaseModel, ValidationError

logger = logging.getLogger(__name__)

_GEMINI_REST_BASE = "https://generativelanguage.googleapis.com/v1beta/models"
_GEMINI_TIMEOUT = 120
_OCR_CONFIDENCE_THRESHOLD = 60
_MIN_TYPED_TEXT_LEN = 100


# ── Pydantic schemas ──────────────────────────────────────────────────────────

class NotesSolutionStepSchema(BaseModel):
    step_type: Literal["context", "given", "formula", "calculation", "result", "insight"]
    content: str
    content_format: str = "text"   # "text" | "table"
    voice_text: Optional[str] = None
    annotation: Optional[dict] = None


class NotesProblemSchema(BaseModel):
    question_text: str
    topic: Optional[str] = None
    given_values: list[str] = []
    find: Optional[str] = None
    method: Optional[str] = None
    difficulty: Literal["easy", "medium", "hard"] = "medium"
    answer: Optional[str] = None
    solution_steps: list[NotesSolutionStepSchema] = []


class NotesExtractionResult(BaseModel):
    problems: list[NotesProblemSchema]


# ── Text extraction ───────────────────────────────────────────────────────────

def extract_notes_text(file_path: str, file_ext: str) -> str:
    ext = file_ext.lower()
    if ext in (".pptx", ".ppt"):
        return _extract_from_pptx(file_path)
    if ext in (".docx", ".doc"):
        return _extract_from_docx(file_path)
    if ext == ".pdf":
        return extract_from_pdf_smart(file_path)
    raise ValueError(f"Unsupported file extension: {file_ext}")


def _extract_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if parts:
                slides_text.append(f"--- Slide {i} ---\n" + "\n".join(parts))
        return "\n\n".join(slides_text)
    except Exception as exc:
        logger.error("[NotesIngestion] PPTX extraction failed: %s", exc)
        return ""


def _extract_from_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as exc:
        logger.error("[NotesIngestion] DOCX extraction failed: %s", exc)
        return ""


def extract_from_pdf_smart(file_path: str) -> str:
    try:
        import fitz
    except ImportError:
        logger.error("[NotesIngestion] PyMuPDF not installed")
        return ""

    doc = fitz.open(file_path)
    pages_text = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if len(text) >= _MIN_TYPED_TEXT_LEN:
            pages_text.append(text)
            logger.debug("[NotesIngestion] Page %d: typed text (%d chars)", page_num, len(text))
        else:
            logger.debug("[NotesIngestion] Page %d: OCR", page_num)
            pages_text.append(ocr_page(page))
    doc.close()
    return "\n\n---\n\n".join(pages_text)


def ocr_page(page) -> str:
    try:
        import fitz
        from PIL import Image
        import pytesseract
        from pytesseract import Output

        img = page.get_pixmap(dpi=200)
        pil_img = Image.frombytes("RGB", [img.width, img.height], img.samples)
        data = pytesseract.image_to_data(pil_img, output_type=Output.DICT)
        confidences = [c for c in data["conf"] if isinstance(c, (int, float)) and c > 0]
        avg_conf = mean(confidences) if confidences else 0.0
        if avg_conf >= _OCR_CONFIDENCE_THRESHOLD:
            return pytesseract.image_to_string(pil_img).strip()
        else:
            return gemini_vision_ocr(pil_img)
    except ImportError:
        try:
            import fitz
            from PIL import Image
            img = page.get_pixmap(dpi=200)
            return gemini_vision_ocr(Image.frombytes("RGB", [img.width, img.height], img.samples))
        except Exception as exc:
            logger.error("[NotesIngestion] OCR fallback failed: %s", exc)
            return "[illegible]"
    except Exception as exc:
        logger.error("[NotesIngestion] OCR page failed: %s", exc)
        return "[illegible]"


def gemini_vision_ocr(pil_image) -> str:
    key = os.getenv("GEMINI_API_KEY", "").strip()
    if not key:
        return "[illegible]"
    buf = io.BytesIO()
    pil_image.save(buf, format="PNG")
    image_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    url = f"{_GEMINI_REST_BASE}/gemini-3.5-flash:generateContent"
    headers = {"x-goog-api-key": key, "Content-Type": "application/json"}
    payload = {
        "contents": [{"parts": [
            {"inlineData": {"mimeType": "image/png", "data": image_b64}},
            {"text": "Extract all text from this image exactly as written. Return only the extracted text."},
        ]}],
        "generationConfig": {"maxOutputTokens": 2048, "temperature": 0.0},
    }
    for attempt, delay in enumerate([0, 2, 4, 8]):
        if delay:
            time.sleep(delay)
        try:
            resp = requests.post(url, headers=headers, json=payload, timeout=_GEMINI_TIMEOUT)
            if resp.status_code == 200:
                return resp.json()["candidates"][0]["content"]["parts"][0]["text"].strip()
            if resp.status_code == 429:
                continue
            logger.error("[NotesIngestion] Gemini Vision HTTP %d", resp.status_code)
        except Exception as exc:
            logger.warning("[NotesIngestion] Gemini Vision attempt %d failed: %s", attempt + 1, exc)
    return "[illegible]"


# ── LLM helpers ───────────────────────────────────────────────────────────────

def _strip_fences(s: str) -> str:
    s = s.strip()
    if s.startswith("```"):
        s = s.split("```", 2)[-1] if s.count("```") >= 2 else s.lstrip("`")
        s = s.lstrip("json").strip()
    if s.endswith("```"):
        s = s[:s.rfind("```")].strip()
    return s


def _call_llm(messages: list, chain: list, max_tokens: int, temperature: float = 0.1) -> str:
    from services.llm_router import chat_with_fallback
    raw, _ = chat_with_fallback(messages=messages, chain=chain, max_tokens=max_tokens, temperature=temperature)
    # Small delay between calls to avoid rate limiting when making multiple sequential calls
    time.sleep(1)
    return _strip_fences(raw)


def _parse_json_array(raw: str) -> list:
    """Parse a JSON array, repairing truncation if needed."""
    raw = raw.strip()
    if not raw.endswith(']'):
        last = raw.rfind('},')
        if last > 0:
            raw = raw[:last + 1] + ']'
        else:
            last = raw.rfind('}')
            if last > 0:
                raw = raw[:last + 1] + ']'
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return []


def _parse_json_object(raw: str) -> dict | None:
    raw = raw.strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return None


# ── Step builder helpers ──────────────────────────────────────────────────────

def _text_step(step_type: str, content: str) -> dict:
    return {"step_type": step_type, "content_format": "text",
            "content": content, "voice_text": content, "annotation": None}


def _table_step(step_type: str, table: dict, voice: str) -> dict:
    return {"step_type": step_type, "content_format": "table",
            "content": json.dumps(table), "voice_text": voice, "annotation": None}


# ── Table generation ──────────────────────────────────────────────────────────

def _gen_table(title: str, table_type: str, instructions: str,
               context: str, chain: list) -> dict | None:
    if table_type == "ledger":
        schema = ('{"type":"ledger","title":"...","left_header":"Dr","right_header":"Cr",'
                  '"rows":[{"left_label":"...","left_amount":"...","right_label":"...","right_amount":"..."}],'
                  '"left_total":"...","right_total":"..."}')
    elif table_type == "balance_sheet":
        schema = ('{"type":"balance_sheet","title":"...","liabilities":[{"label":"...","amount":"..."}],'
                  '"assets":[{"label":"...","amount":"..."}],"total":"..."}')
    else:
        schema = '{"type":"key_value","title":"...","rows":[{"label":"...","value":"...","bold":false}]}'

    prompt = f"""Generate the {title} as a single JSON object.

CONTEXT (source data + pre-computed values):
{context[:5000]}

INSTRUCTIONS:
{instructions}

Schema: {schema}

Rules:
- Use ₹ prefix for all amounts
- Include ALL line items — no placeholders
- For ledger: left_total and right_total MUST be equal (the account must balance)
- Return ONLY the JSON object, no array, no markdown"""

    raw = _call_llm(
        messages=[{"role": "system", "content": "Return only a JSON object."}, {"role": "user", "content": prompt}],
        chain=chain, max_tokens=2000,
    )
    result = _parse_json_object(raw)
    if result:
        logger.info("[NotesIngestion] Generated table: %s (%s rows)", title,
                    len(result.get("rows", result.get("liabilities", []))))
    else:
        logger.warning("[NotesIngestion] Table generation failed: %s", title)
    return result


# ── Accounting step generation ────────────────────────────────────────────────

def _generate_accounting_steps(topic: str, question: str, method: str,
                                answer: str, source_data: str, chain: list) -> list:
    """
    Generate solution steps for accounting problems.
    Strategy:
    1. Pre-compute all adjustment values (separate LLM call → plain text)
    2. Generate text steps (context, given, insight, result)
    3. Generate each table separately using pre-computed values
    """
    steps = []

    # ── Step 1: Pre-compute adjustments ─────────────────────────────────
    pre_prompt = f"""Compute ALL intermediate values needed to solve this accounting problem.
Show each calculation as: "Item = formula = result"

SOURCE: {source_data[:3500]}
QUESTION: {question}
FINAL ANSWER (verify against this): {answer}

Compute: Net Sales, Net Purchases, Gross Profit/Loss, each depreciation,
new bad debts provision, discount on debtors (applied AFTER provision),
interest on drawings, outstanding expense totals, adjusted capital,
net debtors for balance sheet, each fixed asset after depreciation.
Be precise. Show full working for each."""

    precomputed = ""
    try:
        precomputed = _call_llm(
            messages=[{"role": "user", "content": pre_prompt}],
            chain=chain, max_tokens=1200,
        )
        logger.info("[NotesIngestion] Pre-computed %d chars of adjustment values", len(precomputed))
    except Exception as exc:
        logger.warning("[NotesIngestion] Pre-computation failed: %s", exc)

    enriched = f"SOURCE DATA:\n{source_data[:2500]}\n\nPRE-COMPUTED VALUES:\n{precomputed}"

    # ── Step 2: Text steps ───────────────────────────────────────────────
    text_prompt = f"""Generate exactly 4 text steps for this {topic} problem.

QUESTION: {question}
ANSWER: {answer}
PRE-COMPUTED: {precomputed[:1000]}

Return JSON array with these 4 steps:
1. context: what the question requires (1-2 sentences)
2. given: key balances as "Account: ₹Amount; ..." (cover all major items)
3. insight: list each adjustment with its computed value e.g. "Dep. Furniture = 5% × ₹15,500 = ₹775; New Provision = 5% × ₹99,000 = ₹4,950; Discount on Debtors = 2% × ₹94,050 = ₹1,881"
4. result: final figures from answer

Each: {{"step_type":"...","content_format":"text","content":"...","voice_text":"...","annotation":null}}
Max 400 chars per content. Return ONLY the JSON array."""

    try:
        raw = _call_llm(
            messages=[{"role": "system", "content": "Return only a JSON array."},
                      {"role": "user", "content": text_prompt}],
            chain=chain, max_tokens=1500,
        )
        text_steps = _parse_json_array(raw)
        steps.extend(text_steps)
    except Exception as exc:
        logger.warning("[NotesIngestion] Text steps failed: %s", exc)

    # ── Step 3: Tables based on topic ────────────────────────────────────
    topic_lower = topic.lower()
    is_final_accounts = any(k in topic_lower for k in ['trading', 'p&l', 'profit', 'loss', 'final account', 'balance sheet'])
    is_cost_sheet     = 'cost sheet' in topic_lower

    if is_final_accounts:
        # Trading Account
        t = _gen_table(
            title="Trading Account",
            table_type="ledger",
            instructions=f"""Build Trading Account for: {question}
Answer: {answer}

Dr side: To Opening Stock | To Purchases (gross) | Less: Returns Outward → Net Purchases | To Gross Loss c/d (if applicable)
Cr side: By Sales (gross) | Less: Returns Inwards → Net Sales | By Closing Stock | By Gross Profit c/d (if applicable)

Use EXACT amounts from pre-computed values. Both totals MUST balance.""",
            context=enriched,
            chain=chain,
        )
        if t:
            steps.append(_table_step("calculation", t,
                "The Trading Account shows Net Sales vs Cost of Goods to determine Gross Profit or Loss."))

        # P&L Account
        p = _gen_table(
            title="Profit & Loss Account",
            table_type="ledger",
            instructions=f"""Build P&L Account for: {question}
Answer: {answer}

Dr side: To Gross Loss b/d | all indirect expenses with adjustments applied (use pre-computed totals)
  - Rent = base + outstanding
  - Bad Debts = old + further + new provision - old provision
  - Depreciation amounts from pre-computed values
  - Discount on Debtors = 2% of (Debtors - Further BD - New Provision) [NOT of original debtors]
Cr side: By Gross Profit b/d | all incomes | By Interest on Drawings | By Net Loss (balancing)

CRITICAL: Use pre-computed values — do NOT recalculate percentages independently.
Both totals MUST be equal.""",
            context=enriched,
            chain=chain,
        )
        if p:
            steps.append(_table_step("calculation", p,
                "The P&L Account records all indirect incomes and expenses to find Net Profit or Loss."))

        # Balance Sheet
        b = _gen_table(
            title="Balance Sheet",
            table_type="balance_sheet",
            instructions=f"""Build Balance Sheet for: {question}
Answer: {answer}

Liabilities: Capital → Less Net Loss → Less Drawings → Less Interest on Drawings = Adjusted Capital;
  Current Liabilities: Bills Payable, Sundry Creditors, Outstanding expenses

Assets: Fixed Assets (each less its depreciation); Investments;
  Current Assets: Cash in hand, Cash at Bank;
  Sundry Debtors: show gross less bad debts less provision less discount = net;
  Closing Stock

Total Liabilities = Total Assets (use answer for verification)""",
            context=enriched,
            chain=chain,
        )
        if b:
            steps.append(_table_step("calculation", b,
                "The Balance Sheet shows the financial position — assets must equal liabilities plus capital."))

    elif is_cost_sheet:
        c = _gen_table(
            title="Cost Sheet",
            table_type="key_value",
            instructions=f"""Build Cost Sheet for: {question}
Rows in order: Raw Materials Consumed, Direct Wages/Labour, Prime Cost (subtotal, bold),
Factory Overheads itemized, Works/Factory Cost (subtotal, bold),
Office/Admin Overheads itemized, Cost of Production (subtotal, bold),
Selling/Distribution Overheads itemized, Total Cost (bold),
Output Units, Cost per Unit.
Use exact amounts from source. Mark all subtotals bold=true.""",
            context=enriched,
            chain=chain,
        )
        if c:
            steps.append(_table_step("calculation", c,
                "The Cost Sheet systematically builds up from Prime Cost through Works Cost to Total Cost."))

    return steps


def _generate_generic_steps(topic: str, question: str, method: str,
                             answer: str, source_data: str, chain: list) -> list:
    """Generate steps for non-accounting problems (math, NPV, ratios, etc.)."""
    prompt = f"""Generate 7-9 solution steps for this {topic} problem.

QUESTION: {question}
METHOD: {method}
ANSWER: {answer}
SOURCE: {source_data[:3000]}

Return a JSON array. Each step:
{{"step_type":"context|given|formula|calculation|result|insight","content_format":"text","content":"...with actual numbers","voice_text":"...","annotation":null}}
Show real arithmetic. Return ONLY the JSON array."""
    try:
        raw = _call_llm(
            messages=[{"role": "system", "content": "Return only a JSON array."},
                      {"role": "user", "content": prompt}],
            chain=chain, max_tokens=2500,
        )
        return _parse_json_array(raw)
    except Exception:
        return []


# ── Main extraction ───────────────────────────────────────────────────────────

def extract_problems_with_llm(text: str) -> NotesExtractionResult:
    """
    Two-phase extraction:
    1. Extract problem list (question, topic, answer) — one call
    2. Per problem: generate text steps + accounting tables — multiple targeted calls
    """
    from services.llm_config import INGESTION_FALLBACK_CHAIN

    full_text = text[:20000]

    # ── Phase 1: Problem list ────────────────────────────────────────────
    list_prompt = f"""Extract every question/problem from the text below.
For EACH problem return:
- question_text: full question with ALL data values and account names (keep under 800 chars)
- topic: e.g. "Cost Sheet", "Trading Account and Balance Sheet", "NPV"
- given_values: top 6 key values as short strings
- find: what to calculate (one sentence)
- method: approach name
- difficulty: "easy"/"medium"/"hard"
- answer: final answer as a single string with all key figures

Return ONLY: {{"problems": [...]}}. No steps. No markdown.

TEXT:
{full_text[:18000]}"""

    logger.info("[NotesIngestion] Phase 1: extracting problem list")
    try:
        raw1 = _call_llm(
            messages=[{"role": "system", "content": "Return only valid JSON."},
                      {"role": "user", "content": list_prompt}],
            chain=INGESTION_FALLBACK_CHAIN, max_tokens=5000,
        )
        # Repair truncation
        if not raw1.strip().endswith('}'):
            last = raw1.rfind('}')
            if last > 0:
                raw1 = raw1[:last + 1]
                if raw1.count('[') > raw1.count(']'):
                    raw1 += ']}'
        parsed1 = json.loads(raw1)
    except Exception as exc:
        logger.error("[NotesIngestion] Phase 1 failed: %s", exc)
        return NotesExtractionResult(problems=[])

    raw_problems = parsed1.get("problems", []) if isinstance(parsed1, dict) else []
    if not raw_problems:
        return NotesExtractionResult(problems=[])

    logger.info("[NotesIngestion] Phase 1: found %d problems", len(raw_problems))

    # ── Phase 2: Steps per problem ───────────────────────────────────────
    valid_problems: list[NotesProblemSchema] = []

    for i, raw_p in enumerate(raw_problems):
        try:
            # Coerce answer dict → string
            if isinstance(raw_p.get('answer'), dict):
                raw_p = {**raw_p, 'answer': '; '.join(
                    f"{k.replace('_', ' ').title()}: {v}" for k, v in raw_p['answer'].items()
                )}

            question = raw_p.get('question_text', '')
            topic    = raw_p.get('topic', '')
            method   = raw_p.get('method', '')
            answer   = raw_p.get('answer', '') or ''

            topic_lower = topic.lower()
            is_accounting = any(k in topic_lower for k in [
                'trading', 'profit', 'loss', 'balance sheet', 'final account', 'cost sheet', 'p&l'
            ])

            if is_accounting:
                steps_data = _generate_accounting_steps(
                    topic, question, method, answer, full_text, INGESTION_FALLBACK_CHAIN
                )
            else:
                steps_data = _generate_generic_steps(
                    topic, question, method, answer, full_text, INGESTION_FALLBACK_CHAIN
                )

            # Validate each step
            valid_steps = []
            for step in (steps_data if isinstance(steps_data, list) else []):
                try:
                    valid_steps.append(NotesSolutionStepSchema(**step))
                except Exception as e:
                    logger.debug("[NotesIngestion] Invalid step: %s", e)

            raw_p['solution_steps'] = [s.model_dump() for s in valid_steps]
            valid_problems.append(NotesProblemSchema(**raw_p))
            logger.info("[NotesIngestion] Problem %d/%d: %d steps (topic: %s)",
                        i + 1, len(raw_problems), len(valid_steps), topic)

        except Exception as exc:
            logger.warning("[NotesIngestion] Skipping problem %d: %s", i + 1, exc)

    logger.info("[NotesIngestion] Extraction complete: %d problems", len(valid_problems))
    return NotesExtractionResult(problems=valid_problems)


# ── DB persistence ────────────────────────────────────────────────────────────

def save_extraction_result(
    upload_id: str,
    subject_id: str,
    chapter_key: str | None,
    title: str | None,
    result: NotesExtractionResult,
) -> str:
    from db.session import SessionLocal
    from db.models.notes_problem_set import NotesProblemSet
    from db.models.notes_problem import NotesProblem
    from db.models.notes_solution_step import NotesSolutionStep

    problem_set_id = str(uuid.uuid4())
    now = datetime.utcnow()

    with SessionLocal() as session:
        ps = NotesProblemSet(
            id=problem_set_id, upload_id=upload_id, subject_id=subject_id,
            chapter_key=chapter_key, title=title,
            problem_count=len(result.problems), created_at=now,
        )
        session.add(ps)
        session.flush()

        for seq_order, problem in enumerate(result.problems, start=1):
            problem_id = str(uuid.uuid4())
            np = NotesProblem(
                id=problem_id, problem_set_id=problem_set_id,
                sequence_order=seq_order, question_text=problem.question_text,
                topic=problem.topic,
                given_values=json.dumps(problem.given_values) if problem.given_values else None,
                find=problem.find, method=problem.method, difficulty=problem.difficulty,
                answer=problem.answer, created_at=now,
            )
            session.add(np)
            session.flush()

            for step_order, step in enumerate(problem.solution_steps, start=1):
                annotation_json: str | None = None
                if step.annotation and isinstance(step.annotation, dict):
                    target = step.annotation.get("target_text", "")
                    if target and target in problem.question_text:
                        annotation_json = json.dumps(step.annotation)

                ns = NotesSolutionStep(
                    id=str(uuid.uuid4()), problem_id=problem_id,
                    sequence_order=step_order, step_type=step.step_type,
                    content_format=getattr(step, 'content_format', 'text') or 'text',
                    content=step.content, voice_text=step.voice_text,
                    annotation=annotation_json, created_at=now,
                )
                session.add(ns)

        session.commit()

    logger.info("[NotesIngestion] Saved problem_set_id=%s with %d problems",
                problem_set_id, len(result.problems))
    return problem_set_id


# ── Main background task ──────────────────────────────────────────────────────

def run_notes_ingestion(upload_id: str) -> None:
    from db.session import SessionLocal
    from db.models.chapter_upload import ChapterUpload

    logger.info("[NotesIngestion] Starting ingestion for upload_id=%s", upload_id)

    with SessionLocal() as session:
        upload = session.get(ChapterUpload, upload_id)
        if not upload:
            logger.error("[NotesIngestion] Upload %s not found", upload_id)
            return
        file_path = upload.file_path
        file_ext  = Path(upload.original_filename or file_path or "").suffix.lower()
        subject_id  = upload.subject_id
        chapter_key = upload.chapter_key
        title       = upload.chapter_title

    if not file_path or not Path(file_path).exists():
        _mark_failed(upload_id, "Uploaded file not found on disk")
        return

    try:
        text = extract_notes_text(file_path, file_ext)
        if not text.strip():
            _mark_failed(upload_id, "Could not extract any text from the uploaded file")
            return

        result = extract_problems_with_llm(text)
        if not result.problems:
            _mark_failed(upload_id, "No problems could be extracted from this file")
            return

        save_extraction_result(upload_id=upload_id, subject_id=subject_id,
                               chapter_key=chapter_key, title=title, result=result)

        try:
            Path(file_path).unlink(missing_ok=True)
        except Exception as exc:
            logger.warning("[NotesIngestion] Could not delete file: %s", exc)

        _mark_ready(upload_id)
        logger.info("[NotesIngestion] Completed upload_id=%s", upload_id)

    except Exception as exc:
        logger.error("[NotesIngestion] Crashed for upload_id=%s: %s", upload_id, exc, exc_info=True)
        _mark_failed(upload_id, str(exc)[:512])


def _mark_ready(upload_id: str) -> None:
    from db.session import SessionLocal
    from db.models.chapter_upload import ChapterUpload
    with SessionLocal() as session:
        upload = session.get(ChapterUpload, upload_id)
        if upload:
            upload.upload_status = "ready"
            upload.processed_at = datetime.utcnow()
            session.commit()


def _make_friendly_error(raw: str) -> tuple[str, str]:
    """
    Convert a raw technical exception string into:
      - A user-friendly message safe to show in the UI
      - A short error code developers can grep in logs

    The full raw error is always logged separately — never stored in the DB.
    """
    import random
    import string

    # Generate a short 6-char code so devs can find it in logs
    code = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))

    raw_lower = raw.lower()

    # DB / schema errors
    if 'column' in raw_lower and 'does not exist' in raw_lower:
        friendly = "Something went wrong on our end while saving your notes. We're aware and fixing it."
    elif 'relation' in raw_lower and 'does not exist' in raw_lower:
        friendly = "A database table is missing. Please try again in a moment."
    elif 'unique' in raw_lower or 'duplicate key' in raw_lower:
        friendly = "This file appears to have been uploaded already."
    elif 'foreign key' in raw_lower or 'violates' in raw_lower:
        friendly = "Something went wrong while linking your data. Please try again."
    elif 'connection' in raw_lower or 'timeout' in raw_lower or 'refused' in raw_lower:
        friendly = "Couldn't reach the database right now. Please try again in a moment."
    # LLM / AI errors
    elif 'quota' in raw_lower or '429' in raw:
        friendly = "Our AI is a little busy right now. Please try again in a minute."
    elif 'llm' in raw_lower or 'gemini' in raw_lower or 'model' in raw_lower:
        friendly = "The AI had trouble reading your file. Try a text-based PDF or DOCX."
    # File errors
    elif 'not found on disk' in raw_lower or 'no such file' in raw_lower:
        friendly = "The uploaded file couldn't be found. Please upload it again."
    elif 'could not extract' in raw_lower or 'no text' in raw_lower:
        friendly = "Couldn't extract text from this file — it may be a scanned image. Try a text-based PDF."
    elif 'no problems' in raw_lower or 'zero problems' in raw_lower:
        friendly = "No problems were found in this file. Make sure it contains numerical questions."
    # Generic fallback
    else:
        friendly = "Something unexpected happened while processing your notes. Please try again."

    return friendly, code


def _mark_failed(upload_id: str, raw_error: str) -> None:
    from db.session import SessionLocal
    from db.models.chapter_upload import ChapterUpload

    friendly, code = _make_friendly_error(raw_error)

    # Log full technical detail with the code so devs can look it up
    logger.error(
        "[NotesIngestion] FAILED upload_id=%s error_code=%s | %s",
        upload_id, code, raw_error,
    )

    with SessionLocal() as session:
        upload = session.get(ChapterUpload, upload_id)
        if upload:
            upload.upload_status = "failed"
            # Store friendly message + code — never the raw stack trace
            upload.error_message = f"{friendly} (ref: {code})"
            session.commit()
