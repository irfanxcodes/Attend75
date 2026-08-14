"""
Shape Extractor Service

Extracts per-slide shape bounding boxes from a PPTX file.

Called during ingestion (Step 9) BEFORE LibreOffice converts the file
to PDF and BEFORE the original file is deleted.  The results are stored
in lesson_slides.shape_bboxes so the teaching script generator can give
the LLM exact shape coordinates rather than approximate semantic regions.

Output schema per shape:
  {
    "shape_id":     int,    # PPTX internal shape id
    "name":         str,    # shape name from PPTX (e.g. "Title 1", "Content Placeholder 2")
    "type":         str,    # "title" | "body" | "image" | "table" | "other"
    "x":            float,  # left edge, normalised 0.0–1.0 of slide width
    "y":            float,  # top  edge, normalised 0.0–1.0 of slide height
    "w":            float,  # width,     normalised 0.0–1.0
    "h":            float,  # height,    normalised 0.0–1.0
    "text_preview": str,    # first 80 chars of text content (empty for images)
  }

Only PPTX/PPT files are supported.  PDF and DOCX return an empty dict.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_shape_bboxes(file_path: str, file_ext: str) -> dict[int, list[dict]]:
    """
    Return a mapping of slide_number → list-of-shape-dicts.

    Keys are 1-indexed slide numbers matching lesson_slides.slide_number.
    Returns {} for non-PPTX files or on any parsing error (non-fatal).
    """
    ext = file_ext.lower().lstrip(".")
    if ext not in ("pptx", "ppt"):
        return {}

    try:
        return _extract_pptx(file_path)
    except Exception as exc:
        logger.warning("[ShapeExtractor] Failed for %s: %s", file_path, exc)
        return {}


def _extract_pptx(file_path: str) -> dict[int, list[dict]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    slide_w = prs.slide_width.emu or 9_144_000
    slide_h = prs.slide_height.emu or 6_858_000

    result: dict[int, list[dict]] = {}

    for slide_no, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            # Skip invisible / zero-size shapes
            if not shape.width or not shape.height:
                continue

            shape_type = _classify_shape(shape)

            # Clamp and normalise
            left   = max(0, shape.left   or 0)
            top    = max(0, shape.top    or 0)
            width  = max(1, shape.width  or 1)
            height = max(1, shape.height or 1)

            x = round(left   / slide_w, 4)
            y = round(top    / slide_h, 4)
            w = round(width  / slide_w, 4)
            h = round(height / slide_h, 4)

            # Cap at slide boundary
            x = min(x, 1.0)
            y = min(y, 1.0)
            w = min(w, 1.0 - x)
            h = min(h, 1.0 - y)

            # Text preview (title/body shapes only)
            text_preview = ""
            if shape.has_text_frame:
                text_preview = shape.text_frame.text[:80].strip()

            shapes.append({
                "shape_id":     shape.shape_id,
                "name":         shape.name or "",
                "type":         shape_type,
                "x":            x,
                "y":            y,
                "w":            w,
                "h":            h,
                "text_preview": text_preview,
            })

        if shapes:
            result[slide_no] = shapes

    logger.debug("[ShapeExtractor] Extracted bboxes for %d slides from %s",
                 len(result), file_path)
    return result


def _classify_shape(shape) -> str:
    """
    Classify a PPTX shape into one of: title | body | image | table | other.
    Uses placeholder index first (most reliable), then shape type.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    ph = getattr(shape, "placeholder_format", None)
    if ph is not None:
        idx = ph.idx
        if idx in (0, 1):
            return "title"
        # idx 2+ are typically body/content placeholders
        return "body"

    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        if st == MSO_SHAPE_TYPE.TABLE:
            return "table"
    except Exception:
        pass

    if shape.has_text_frame and shape.text_frame.text.strip():
        return "body"

    return "other"
