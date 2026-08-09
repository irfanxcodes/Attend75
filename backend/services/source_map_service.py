"""
Source Map Service — StudyMe 2.0 Phase 3

Extracts a slide-by-slide / page-by-page structural map from the original
uploaded document. Used to power the Source (PPT) viewer in WorkspacePlayer.

Each entry in the map represents one slide or page with:
  - number (1-indexed)
  - title (first heading or title shape)
  - body_preview (first 300 chars of body text)
  - full_text (complete text for search/highlighting)
  - element_type ('slide' for PPTX, 'page' for PDF)

The original file is deleted after ingestion (by design), so:
  - If file still exists on disk: parse it fresh
  - If file is gone: reconstruct from ai_concepts source_page data
    (good enough for concept navigation even without the original file)
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def build_source_map(upload_id: str) -> list[dict]:
    """
    Build a source map for a chapter upload.

    Returns list of:
      {number, title, body_preview, full_text, element_type}

    Falls back gracefully if file is missing.
    """
    from db.session import SessionLocal
    from db.models.chapter_upload import ChapterUpload
    from db.models.ai_concept import AIConcept

    with SessionLocal() as session:
        upload = session.get(ChapterUpload, upload_id)
        if not upload:
            return []

        file_path = upload.file_path
        original_filename = upload.original_filename or ""

        # Determine file type from original filename
        ext = Path(original_filename).suffix.lower() if original_filename else ""

        # Try parsing from disk first (file may still be there in dev)
        if file_path and Path(file_path).exists():
            try:
                if ext in (".pptx", ".ppt"):
                    result = _parse_pptx_map(file_path)
                    if result:
                        logger.info("[SourceMap] Built PPTX map: %d slides", len(result))
                        return result
                elif ext == ".pdf":
                    result = _parse_pdf_map(file_path)
                    if result:
                        logger.info("[SourceMap] Built PDF map: %d pages", len(result))
                        return result
            except Exception as exc:
                logger.warning("[SourceMap] File parse failed: %s", exc)

        # Fallback: reconstruct from ai_concepts source_page metadata
        concepts = (
            session.query(AIConcept)
            .filter(AIConcept.upload_id == upload_id)
            .order_by(AIConcept.source_page, AIConcept.sequence_order)
            .all()
        )

        if not concepts:
            return []

        logger.info("[SourceMap] Reconstructing from %d concepts (file unavailable)", len(concepts))
        return _reconstruct_from_concepts(concepts, ext)


def _parse_pptx_map(file_path: str) -> list[dict]:
    """Extract slide structure from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_parts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            ph = getattr(shape, 'placeholder_format', None)
            is_title = ph is not None and ph.idx in (0, 1)

            if is_title and not title_text:
                title_text = text[:120]
            else:
                body_parts.append(text)

        full_text = "\n".join([title_text] + body_parts) if title_text else "\n".join(body_parts)
        body_preview = "\n".join(body_parts)[:300]

        slides.append({
            "number": i,
            "title": title_text or f"Slide {i}",
            "body_preview": body_preview,
            "full_text": full_text,
            "element_type": "slide",
        })

    return slides


def _parse_pdf_map(file_path: str) -> list[dict]:
    """Extract page structure from a PDF file."""
    from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
    from docling.datamodel.document import InputDocument
    from docling.datamodel.base_models import InputFormat

    path = Path(file_path)
    doc_in = InputDocument(
        path_or_stream=path,
        format=InputFormat.PDF,
        backend=PyPdfiumDocumentBackend,
    )
    backend = doc_in._backend
    pages = []

    for pg in backend.iter_pages():
        cells = list(pg.get_text_cells())
        lines = [c.text.strip() for c in cells if c.text.strip()]
        if not lines:
            continue

        title = lines[0][:120] if lines else f"Page {pg.page_no + 1}"
        body_lines = lines[1:] if len(lines) > 1 else []
        body_preview = " ".join(body_lines)[:300]
        full_text = " ".join(lines)

        pages.append({
            "number": pg.page_no + 1,
            "title": title,
            "body_preview": body_preview,
            "full_text": full_text,
            "element_type": "page",
        })

    backend.unload()
    return pages


def _reconstruct_from_concepts(concepts, ext: str) -> list[dict]:
    """
    Reconstruct a minimal source map from concept source_page metadata.
    Used when the original file has been deleted.
    Groups concepts by source_page to create one entry per unique page/slide.
    """
    element_type = "slide" if ext in (".pptx", ".ppt") else "page"

    # Group concepts by page
    page_map: dict[int, list] = {}
    for concept in concepts:
        page = concept.source_page or 0
        if page not in page_map:
            page_map[page] = []
        page_map[page].append(concept)

    entries = []
    for page_no in sorted(page_map.keys()):
        if page_no == 0:
            continue  # Skip unpositioned concepts
        page_concepts = page_map[page_no]
        # Use heading of first concept on this page as title
        title = page_concepts[0].source_heading or page_concepts[0].title
        # Body: concept titles on this page
        body_preview = "\n".join(c.title for c in page_concepts[:4])
        full_text = " | ".join(c.title for c in page_concepts)

        entries.append({
            "number": page_no,
            "title": title[:120],
            "body_preview": body_preview,
            "full_text": full_text,
            "element_type": element_type,
            "concepts": [c.title for c in page_concepts],  # bonus metadata
        })

    return entries
