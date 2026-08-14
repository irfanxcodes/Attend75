"""
Slide Renderer Service

Converts uploaded PPT/PDF files into per-slide WebP images.

Strategy (in priority order):
  1. LibreOffice headless → PDF → PyMuPDF renders each page
     ✓ Pixel-perfect: preserves original fonts, colors, images, layouts
     ✓ Works for PPTX, PPT, PDF, DOCX
     ✓ Available on Ubuntu via: apt install libreoffice-common libreoffice-impress

  2. PyMuPDF direct (PDF only fallback)
     ✓ Used when LibreOffice is not installed (dev machines)
     ✓ Still high quality for native PDFs

  3. python-pptx + Pillow (last resort PPTX fallback)
     ✗ No original fonts/colors/images — plain reconstruction
     ✗ Only used when both LibreOffice and direct PDF path are unavailable

Output: WebP images at 1280px wide, quality 82.
Storage: delegated to slide_storage.py (local disk in dev, R2 in prod).
Deduplication: caller checks slide_image_exists() before calling render_slides().
"""

import io
import logging
import os
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

_TARGET_W   = 1280   # output width in pixels
_WEBP_Q     = 82     # WebP quality (82 is visually sharp, ~25-60 KB/slide)
_LO_TIMEOUT = 120    # seconds to wait for LibreOffice conversion


# ── Public API ────────────────────────────────────────────────────────────────

def render_slides(upload_id: str, file_path: str, file_ext: str) -> list[dict]:
    """
    Render all slides/pages of a file to WebP images and save via slide_storage.

    Returns list of dicts:
      [{"slide_number": 1, "url": "...", "width_px": 1280, "height_px": 720}, ...]

    Returns [] on complete failure (non-fatal — lesson still works without slides).
    Skips slides that already exist in storage (idempotent — retries are free).
    Respects MAX_SLIDES_PER_UPLOAD cap AND all three storage guards in
    storage_cap_service.  Raises StorageCapExceeded if any guard is hit —
    callers must treat this as a hard error (not a silent empty return).

    Fail-closed: if the cap check itself fails (DB down etc.), StorageCapExceeded
    is raised and the upload is blocked.
    """
    from services.slide_storage import MAX_SLIDES_PER_UPLOAD, save_slide_image, slide_image_exists
    from services.storage_cap_service import StorageCapExceeded, check_and_reserve, release_reservation

    ext = file_ext.lower().lstrip(".")
    results = []

    try:
        # ── Choose rendering path ─────────────────────────────────────────
        if _libreoffice_available() and ext in ("pptx", "ppt", "pdf", "docx", "doc"):
            slide_images = _render_via_libreoffice(file_path, ext)
        elif ext == "pdf":
            slide_images = _render_pdf_direct(file_path)
        elif ext in ("pptx", "ppt"):
            slide_images = _render_pptx_pillow(file_path)
        else:
            slide_images = _render_pdf_direct(file_path)

        if not slide_images:
            logger.warning("[SlideRenderer] No slides extracted from %s", file_path)
            return []

        # ── Per-upload slide count cap ────────────────────────────────────
        if len(slide_images) > MAX_SLIDES_PER_UPLOAD:
            logger.warning("[SlideRenderer] Capping %d slides to %d for upload %s",
                           len(slide_images), MAX_SLIDES_PER_UPLOAD, upload_id)
            slide_images = slide_images[:MAX_SLIDES_PER_UPLOAD]

        # ── Identify genuinely new slides (dedup — retries don't re-charge) ──
        # Only slides that do NOT already exist in storage consume cap budget.
        new_slide_indices: list[int] = []
        for i, (img_bytes, w, h) in enumerate(slide_images):
            if not slide_image_exists(upload_id, i + 1):
                new_slide_indices.append(i)

        n_new = len(new_slide_indices)

        # ── Atomic cap reservation (all three guards, fail-closed) ────────
        # actual_bytes is the REAL byte size of each new WebP — no estimates.
        if n_new > 0:
            actual_bytes = [len(slide_images[i][0]) for i in new_slide_indices]
            check_and_reserve(n_new, actual_bytes)  # raises StorageCapExceeded if any guard fails

        # ── Write to storage (release reservation on partial failure) ─────
        stored_new = 0
        stored_bytes: list[int] = []
        try:
            for i, (img_bytes, w, h) in enumerate(slide_images):
                slide_no = i + 1

                if i not in new_slide_indices:
                    # Already stored — return existing URL, no R2 PUT needed.
                    from services.slide_storage import get_slide_image_url
                    url = get_slide_image_url(upload_id, slide_no) or ""
                    results.append({"slide_number": slide_no, "url": url, "width_px": w, "height_px": h})
                    continue

                url = save_slide_image(upload_id, slide_no, img_bytes)
                results.append({"slide_number": slide_no, "url": url, "width_px": w, "height_px": h})
                stored_new += 1
                stored_bytes.append(len(img_bytes))

        except Exception as store_exc:
            # Release reservation for slides we reserved but did NOT write.
            not_written = n_new - stored_new
            if not_written > 0:
                unreserved_bytes = actual_bytes[stored_new:]  # bytes we never PUT
                release_reservation(not_written, unreserved_bytes)
                logger.warning(
                    "[SlideRenderer] Released reservation for %d unwritten slides (%d bytes)",
                    not_written, sum(unreserved_bytes),
                )
            logger.error(
                "[SlideRenderer] Storage failed mid-upload for %s after %d/%d slides: %s",
                upload_id, stored_new, n_new, store_exc, exc_info=True,
            )
            raise

        logger.info(
            "[SlideRenderer] Rendered %d slides (%d new, %d bytes) for upload_id=%s",
            len(results), stored_new, sum(stored_bytes), upload_id,
        )
        return results

    except StorageCapExceeded:
        raise  # hard error — do not swallow
    except Exception as exc:
        logger.error("[SlideRenderer] Render failed for upload_id=%s: %s", upload_id, exc, exc_info=True)
        return []


# ── Strategy 1: LibreOffice → PDF → PyMuPDF ──────────────────────────────────

def _libreoffice_available() -> bool:
    """Check if LibreOffice is installed on this machine."""
    return _libreoffice_cmd() is not None


def _libreoffice_cmd() -> str | None:
    """Find the LibreOffice executable path."""
    candidates = [
        "libreoffice",
        "soffice",
        # macOS application bundle (direct download install)
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        # macOS Homebrew cask
        "/opt/homebrew/bin/libreoffice",
        "/usr/local/bin/libreoffice",
        # Linux
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
    ]
    for cmd in candidates:
        try:
            r = subprocess.run([cmd, "--version"], capture_output=True, timeout=10)
            if r.returncode == 0:
                return cmd
        except (FileNotFoundError, subprocess.TimeoutExpired, PermissionError):
            continue
    return None


def _render_via_libreoffice(file_path: str, ext: str) -> list[tuple[bytes, int, int]]:
    """
    Use LibreOffice headless to convert the file to PDF, then render
    each page to a WebP image via PyMuPDF.

    This preserves: fonts, colors, images, gradients, shadows, tables, SmartArt.
    """
    import fitz
    from PIL import Image

    with tempfile.TemporaryDirectory() as tmpdir:
        # For PDF input, skip LibreOffice conversion step
        if ext == "pdf":
            pdf_path = file_path
        else:
            pdf_path = _libreoffice_convert_to_pdf(file_path, tmpdir)
            if not pdf_path:
                logger.warning("[SlideRenderer] LibreOffice conversion failed, falling back")
                return []

        return _pdf_to_images(pdf_path)


def _libreoffice_convert_to_pdf(file_path: str, output_dir: str) -> str | None:
    """Convert a PPTX/DOCX/etc to PDF using LibreOffice headless."""
    lo = _libreoffice_cmd()
    try:
        result = subprocess.run(
            [
                lo,
                "--headless",
                "--norestore",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                file_path,
            ],
            capture_output=True,
            text=True,
            timeout=_LO_TIMEOUT,
        )
        if result.returncode != 0:
            logger.error("[SlideRenderer] LibreOffice stderr: %s", result.stderr[:500])
            return None

        # LibreOffice outputs {basename}.pdf in outdir
        stem = Path(file_path).stem
        pdf_path = Path(output_dir) / f"{stem}.pdf"
        if pdf_path.exists():
            return str(pdf_path)

        # Sometimes LibreOffice changes the name — find any PDF in outdir
        pdfs = list(Path(output_dir).glob("*.pdf"))
        return str(pdfs[0]) if pdfs else None

    except subprocess.TimeoutExpired:
        logger.error("[SlideRenderer] LibreOffice timed out after %ds", _LO_TIMEOUT)
        return None
    except Exception as exc:
        logger.error("[SlideRenderer] LibreOffice error: %s", exc)
        return None


# ── Strategy 2: PyMuPDF direct (PDF) ─────────────────────────────────────────

def _render_pdf_direct(file_path: str) -> list[tuple[bytes, int, int]]:
    """Render a PDF directly using PyMuPDF — high quality for native PDFs."""
    return _pdf_to_images(file_path)


def _pdf_to_images(pdf_path: str) -> list[tuple[bytes, int, int]]:
    """Render all pages of a PDF to WebP images via PyMuPDF."""
    import fitz
    from PIL import Image

    results = []
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            page_w = page.rect.width
            scale = _TARGET_W / page_w if page_w > 0 else 1.5
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            buf = io.BytesIO()
            img.save(buf, format="WEBP", quality=_WEBP_Q, method=4)
            results.append((buf.getvalue(), pix.width, pix.height))

        doc.close()
    except Exception as exc:
        logger.error("[SlideRenderer] PDF render error: %s", exc)

    return results


# ── Strategy 3: python-pptx + Pillow (last resort) ───────────────────────────

def _render_pptx_pillow(file_path: str) -> list[tuple[bytes, int, int]]:
    """
    Fallback PPTX renderer using python-pptx shape extraction + Pillow drawing.
    No LibreOffice needed, but output is a plain reconstruction without original
    fonts, colors, or images. Used only when LibreOffice is not installed.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageDraw

    _SLIDE_W = 1280
    _SLIDE_H = 720

    try:
        prs = Presentation(file_path)
        slide_emu_w = prs.slide_width.emu or 9144000
        slide_emu_h = prs.slide_height.emu or 6858000
        scale_x = _SLIDE_W / slide_emu_w
        scale_y = _SLIDE_H / slide_emu_h

        results = []
        for slide in prs.slides:
            img = _render_pptx_slide(slide, scale_x, scale_y, _SLIDE_W, _SLIDE_H)
            buf = io.BytesIO()
            img.save(buf, format="WEBP", quality=_WEBP_Q, method=4)
            results.append((buf.getvalue(), _SLIDE_W, _SLIDE_H))

        return results
    except Exception as exc:
        logger.error("[SlideRenderer] Pillow PPTX fallback error: %s", exc)
        return []


def _render_pptx_slide(slide, scale_x, scale_y, W, H):
    from PIL import Image, ImageDraw
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    bg = _pptx_slide_bg(slide)
    draw.rectangle([0, 0, W, H], fill=bg)
    is_dark = sum(bg) < 400
    draw.rectangle([0, 0, W, 6], fill=(40, 50, 120) if not is_dark else (255, 255, 255))

    for shape in slide.shapes:
        left   = max(0, min(int((shape.left   or 0) * scale_x), W - 1))
        top    = max(0, min(int((shape.top    or 0) * scale_y), H - 1))
        width  = max(1, min(int((shape.width  or W) * scale_x), W - left))
        height = max(1, min(int((shape.height or H) * scale_y), H - top))

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic_data = shape.image.blob
                from PIL import Image as PILImage
                pic = PILImage.open(io.BytesIO(pic_data)).convert("RGB")
                pic = pic.resize((width, height), PILImage.LANCZOS)
                img.paste(pic, (left, top))
            except Exception:
                pass
            continue

        if shape.has_table:
            _pptx_draw_table(draw, shape.table, left, top, width, height, is_dark)
            continue

        if shape.has_text_frame:
            _pptx_draw_text(draw, shape, left, top, width, height, is_dark)

    return img


def _pptx_slide_bg(slide):
    try:
        fc = slide.background.fill.fore_color
        if fc.rgb:
            return (fc.rgb.red, fc.rgb.green, fc.rgb.blue)
    except Exception:
        pass
    return (255, 255, 255)


def _pptx_draw_text(draw, shape, left, top, width, height, is_dark):
    tf = shape.text_frame
    if not tf.text.strip():
        return
    ph = getattr(shape, "placeholder_format", None)
    is_title = ph is not None and ph.idx in (0, 1)
    tc = (20, 20, 60) if not is_dark else (240, 240, 255)
    bc = (50, 50, 80) if not is_dark else (200, 210, 240)
    font_t = _pf(bold=True, size=34)
    font_b = _pf(bold=False, size=19)
    pad = 12
    x, y = left + pad, top + pad
    mw, my = width - pad * 2, top + height - pad

    for para in tf.paragraphs:
        text = para.text.strip()
        if not text or y >= my:
            continue
        font = font_t if is_title else font_b
        color = tc if is_title else bc
        if not is_title:
            draw.ellipse([x, y + 7, x + 7, y + 14], fill=(100, 130, 220))
        tx = x + (0 if is_title else 16)
        for line in _wrap(draw, text, font, mw - (0 if is_title else 16))[:4]:
            if y >= my:
                break
            draw.text((tx, y), line, fill=color, font=font)
            y += _lh(font)
        y += 4


def _pptx_draw_table(draw, table, left, top, width, height, is_dark):
    nc, nr = len(table.columns), len(table.rows)
    if not nc or not nr:
        return
    cw = width // nc
    rh = min(height // nr, 36)
    hbg = (50, 80, 160)
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cx, cy = left + c * cw, top + r * rh
            bg = hbg if r == 0 else ((240, 244, 255) if r % 2 else (250, 252, 255))
            draw.rectangle([cx, cy, cx + cw - 1, cy + rh - 1], fill=bg,
                           outline=(200, 210, 230))
            txt = cell.text.strip()[:30]
            if txt:
                draw.text((cx + 6, cy + 6), txt,
                          fill=(255, 255, 255) if r == 0 else (30, 40, 80),
                          font=_pf(bold=(r == 0), size=13))


_FONT_CACHE: dict = {}

def _pf(bold=False, size=18):
    from PIL import ImageFont
    key = f"{'b' if bold else 'r'}{size}"
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    candidates = (
        ["/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
         "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
         "/System/Library/Fonts/Helvetica.ttc"]
        if bold else
        ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
         "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
         "/System/Library/Fonts/Helvetica.ttc"]
    )
    font = None
    for p in candidates:
        try:
            font = ImageFont.truetype(p, size)
            break
        except (IOError, OSError):
            continue
    if not font:
        font = ImageFont.load_default()
    _FONT_CACHE[key] = font
    return font


def _wrap(draw, text, font, max_w):
    words, lines, cur = text.split(), [], ""
    for w in words:
        test = (cur + " " + w).strip()
        try:
            bx = draw.textbbox((0, 0), test, font=font)
            tw = bx[2] - bx[0]
        except Exception:
            tw = len(test) * 10
        if tw <= max_w:
            cur = test
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


_DUMMY = None
def _lh(font):
    global _DUMMY
    try:
        from PIL import Image, ImageDraw
        if _DUMMY is None:
            _DUMMY = Image.new("RGB", (1, 1))
        bx = ImageDraw.Draw(_DUMMY).textbbox((0, 0), "Ay", font=font)
        return int((bx[3] - bx[1]) * 1.35)
    except Exception:
        return 26
